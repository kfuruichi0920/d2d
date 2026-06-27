"""PowerPoint (.pptx) 抽出コマンド"""

from __future__ import annotations
import os
import uuid
from typing import Callable


def run(params: dict, progress: Callable[[str], None]) -> dict:
    try:
        from pptx import Presentation
        from pptx.util import Pt
    except ImportError:
        raise RuntimeError("python-pptx is not installed. Run: pip install python-pptx")

    blob_path: str = params["blob_path"]
    progress(f"PowerPoint 文書を読み込み中: {os.path.basename(blob_path)}")

    prs = Presentation(blob_path)
    items = []
    structure = []

    for slide_idx, slide in enumerate(prs.slides):
        slide_no = slide_idx + 1
        progress(f"  スライド {slide_no} / {len(prs.slides)}")

        slide_uid = str(uuid.uuid4())
        slide_texts = []

        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = "".join(run.text for run in para.runs).strip()
                    if text:
                        slide_texts.append(text)

        items.append({
            "item_uid": slide_uid,
            "item_type": "text",
            "slide_number": slide_no,
            "texts": slide_texts,
            "text": "\n".join(slide_texts),
        })
        structure.append({
            "uid": slide_uid,
            "title": slide_texts[0] if slide_texts else f"スライド {slide_no}",
            "children": [],
        })

    progress(f"抽出完了: {len(items)} スライド")
    return {"items": items, "structure_json": structure}
